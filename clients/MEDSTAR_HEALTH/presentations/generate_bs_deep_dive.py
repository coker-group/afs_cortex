"""
MedStar Health — Balance Sheet Deep Dive PPTX Generator
Builds a 15-slide presentation from the spec using python-pptx.
No external template required — creates a blank widescreen deck
and applies Acuvance Coker styling programmatically.

Usage (from workspace root):
    pip install python-pptx
    python clients/MEDSTAR_HEALTH/presentations/generate_bs_deep_dive.py

Output:
    clients/MEDSTAR_HEALTH/presentations/deep_dive_balance_sheet.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.oxml.ns import qn
import os

AC_BLUE  = RGBColor(0x00, 0x4B, 0xFF)
AC_DEEP  = RGBColor(0x0C, 0x36, 0xB4)
AC_MID   = RGBColor(0x4A, 0x78, 0xFF)
AC_PALE  = RGBColor(0x7B, 0xA3, 0xFF)
AC_GHOST = RGBColor(0xB3, 0xC8, 0xFF)
AC_TEAL  = RGBColor(0x4C, 0xC6, 0xC6)
AC_AMBER = RGBColor(0xF4, 0xB7, 0x3F)
AC_GREEN = RGBColor(0x16, 0x98, 0x73)
AC_RED   = RGBColor(0xB8, 0x32, 0x32)
AC_DARK  = RGBColor(0x22, 0x27, 0x2D)
AC_GRAY  = RGBColor(0x7A, 0x7D, 0x81)
AC_LGRAY = RGBColor(0xF3, 0xF4, 0xF5)
AC_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
AC_GRID  = RGBColor(0xE8, 0xE9, 0xEA)
AC_ALT   = RGBColor(0xCC, 0xDB, 0xFF)
AC_BORDER = RGBColor(0x0C, 0x36, 0xB4)

YEARS = ['FY2018','FY2019','FY2020','FY2021','FY2022','FY2023','FY2024','FY2025']
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN_L = Inches(0.8)
MARGIN_R = Inches(0.8)
CONTENT_W = Inches(11.733)
FONT = 'Arial'

OUT_DIR = os.path.dirname(os.path.abspath(__file__))
OUT_PATH = os.path.join(OUT_DIR, 'deep_dive_balance_sheet.pptx')


def _set_text(tf, text, size=10, bold=False, color=AC_DARK, align=PP_ALIGN.LEFT):
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _add_run(p, text, size=10, bold=False, color=AC_DARK):
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return run


def _add_textbox(slide, left, top, width, height, text,
                 size=10, bold=False, color=AC_DARK, align=PP_ALIGN.LEFT,
                 anchor=MSO_ANCHOR.TOP, fill_color=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    if anchor:
        tf.paragraphs[0].alignment = align
        try:
            txBox.text_frame._txBody.attrib[qn('anchor')] = {
                MSO_ANCHOR.TOP: 't', MSO_ANCHOR.MIDDLE: 'ctr', MSO_ANCHOR.BOTTOM: 'b'
            }.get(anchor, 't')
        except Exception:
            pass
    if fill_color:
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = fill_color
    _set_text(tf, text, size, bold, color, align)
    return txBox


def _add_rect(slide, left, top, width, height, fill_color, border_color=None, corner=Inches(0.06)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    if corner:
        try:
            shape.adjustments[0] = 0.04
        except Exception:
            pass
    return shape


def _add_blue_rule(slide):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Pt(4))
    line.fill.solid()
    line.fill.fore_color.rgb = AC_BLUE
    line.line.fill.background()


def _add_title(slide, text):
    _add_blue_rule(slide)
    _add_textbox(slide, MARGIN_L, Inches(0.25), CONTENT_W, Inches(0.7),
                 text, size=22, bold=True, color=AC_BLUE)


def _add_source(slide, text):
    _add_textbox(slide, MARGIN_L, Inches(6.95), CONTENT_W, Inches(0.3),
                 text, size=7, color=AC_GRAY)


def _add_footer(slide, page_num, total=15):
    _add_textbox(slide, MARGIN_L, Inches(7.15), Inches(2), Inches(0.25),
                 'www.cokergroup.com', size=7, color=AC_GRAY)
    _add_textbox(slide, Inches(5.5), Inches(7.15), Inches(3), Inches(0.25),
                 'MedStar Health \u2014 Balance Sheet Deep Dive', size=7, color=AC_GRAY,
                 align=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(11), Inches(7.15), Inches(1.5), Inches(0.25),
                 f'pg. {page_num} / {total}', size=7, color=AC_GRAY, align=PP_ALIGN.RIGHT)


def _add_callout_box(slide, left, top, width, height, title, body, border_color=AC_BLUE):
    bg = _add_rect(slide, left, top, width, height, AC_LGRAY)
    bdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Pt(4), height)
    bdr.fill.solid()
    bdr.fill.fore_color.rgb = border_color
    bdr.line.fill.background()
    tb = slide.shapes.add_textbox(left + Pt(14), top + Pt(8), width - Pt(22), height - Pt(16))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    _add_run(p, title, size=10, bold=True, color=AC_DARK)
    p2 = tf.add_paragraph()
    p2.space_before = Pt(3)
    _add_run(p2, body, size=8, color=AC_GRAY)


def _style_chart_axes(chart, y_fmt='$#,##0"M"', y_min=None, y_max=None,
                      show_x_grid=False, show_y_grid=True):
    cat_ax = chart.category_axis
    cat_ax.has_major_gridlines = show_x_grid
    cat_ax.tick_labels.font.name = FONT
    cat_ax.tick_labels.font.size = Pt(9)
    cat_ax.tick_labels.font.color.rgb = AC_DARK
    cat_ax.format.line.fill.background()
    cat_ax.major_tick_mark = 2

    val_ax = chart.value_axis
    val_ax.has_major_gridlines = show_y_grid
    if show_y_grid:
        val_ax.major_gridlines.format.line.color.rgb = AC_GRID
        val_ax.major_gridlines.format.line.width = Pt(0.5)
    if y_fmt:
        val_ax.tick_labels.number_format = y_fmt
    val_ax.tick_labels.font.name = FONT
    val_ax.tick_labels.font.size = Pt(9)
    val_ax.tick_labels.font.color.rgb = AC_DARK
    val_ax.format.line.fill.background()
    if y_min is not None:
        val_ax.minimum_scale = y_min
    if y_max is not None:
        val_ax.maximum_scale = y_max


def _color_series(plot, idx, color, fill_opacity=None):
    series = plot.series[idx]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = color
    if fill_opacity is not None:
        try:
            solidFill = series.format.fill._fill
            srgb = solidFill.find(qn('a:srgbClr'))
            if srgb is not None:
                alpha = srgb.makeelement(qn('a:alpha'), {})
                alpha.set('val', str(int(fill_opacity * 1000)))
                srgb.append(alpha)
        except Exception:
            pass


def build_deck():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    # ── SLIDE 1: COVER ──
    s1 = prs.slides.add_slide(blank)
    bg_shape = _add_rect(s1, 0, 0, Inches(6.5), SLIDE_H, AC_BLUE)
    _add_rect(s1, Inches(6.5), 0, Inches(6.833), SLIDE_H, AC_DEEP)
    _add_rect(s1, Inches(5.8), Inches(1.0), Inches(1.8), Inches(5.5), AC_MID)
    _add_rect(s1, Inches(6.15), Inches(0.5), Inches(1.8), Inches(6.5), AC_PALE)

    _add_textbox(s1, Inches(0.8), Inches(0.5), Inches(3), Inches(0.4),
                 'Balance Sheet & Credit', size=9, bold=True, color=AC_DARK,
                 fill_color=AC_WHITE)
    _add_textbox(s1, Inches(0.8), Inches(2.5), Inches(5.0), Inches(1.2),
                 'MedStar Health, Inc.', size=36, bold=True, color=AC_WHITE)
    _add_textbox(s1, Inches(0.8), Inches(3.8), Inches(5.0), Inches(0.8),
                 'Balance Sheet Health & Trajectory \u2014 Deep Dive Analysis',
                 size=16, color=AC_WHITE)
    _add_textbox(s1, Inches(0.8), Inches(6.7), Inches(5.0), Inches(0.4),
                 'CONFIDENTIAL  |  May 2026  |  Domain Deep Dive',
                 size=9, color=AC_WHITE)

    logo_path = os.path.join(os.path.dirname(OUT_DIR), '..', '..',
                             'acuvance_coker_design_system', 'assets',
                             'acuvance-coker-logo-white.png')
    if os.path.exists(logo_path):
        s1.shapes.add_picture(logo_path, Inches(0.8), Inches(6.2), Inches(2.0))

    # ── SLIDE 2: BS SNAPSHOT ──
    s2 = prs.slides.add_slide(blank)
    _add_title(s2, 'Balance Sheet Snapshot: Structural Strength Masking Operational Fragility')

    kpis = [
        ('$4.3B', 'Net Assets (Unrest.)', AC_WHITE),
        ('28.5%', 'Debt / Cap', AC_GREEN),
        ('122d', 'Days Cash', AC_RED),
        ('51d', 'Days A/R', AC_AMBER),
        ('$14M', 'Pension Liab.', AC_GREEN),
    ]
    kpi_w = Inches(2.15)
    kpi_h = Inches(0.85)
    for i, (val, lbl, val_color) in enumerate(kpis):
        left = MARGIN_L + Inches(i * 2.3)
        top = Inches(1.15)
        box = _add_rect(s2, left, top, kpi_w, kpi_h, AC_DARK)
        tb = s2.shapes.add_textbox(left + Pt(6), top + Pt(6), kpi_w - Pt(12), kpi_h - Pt(12))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _add_run(p, val, size=24, bold=True, color=val_color)
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        _add_run(p2, lbl, size=8, color=AC_WHITE)

    callouts = [
        (AC_GREEN, 'The Good: Capital Structure Is the Strongest in a Decade',
         'Unrestricted net assets grew 140% since FY2018 ($1.8B\u2192$4.3B). Debt-to-cap at 28.5% is well below the 33% A-rated ceiling. Pension virtually eliminated ($629M\u2192$14M).'),
        (AC_RED, 'The Concern: Liquidity and Cash Generation Lag',
         'Days cash stuck at ~122 vs. 225-day A-rated median. Operating CF margin at 2.0% vs. 7% benchmark. CapEx/depreciation at 8% = severe deferred maintenance.'),
        (AC_AMBER, 'The Risk: Growth Driven by Investment Returns, Not Operations',
         'Cumulative operating income FY2018\u2013FY2025: $1.3B. But net assets grew $2.5B. The difference = investment market gains.'),
        (AC_BLUE, 'The Opportunity: Debt Capacity Headroom',
         'At 28.5% debt-to-cap vs. 33% ceiling, MedStar has $300\u2013400M in unused strategic debt capacity.'),
    ]
    co_w = Inches(5.6)
    co_h = Inches(1.15)
    for i, (bc, title, body) in enumerate(callouts):
        col = i % 2
        row = i // 2
        left = MARGIN_L + Inches(col * 5.9)
        top = Inches(2.25) + Inches(row * 1.3)
        _add_callout_box(s2, left, top, co_w, co_h, title, body, bc)

    _add_source(s2, 'Source: Audited Financial Statements, FY2018\u20132025 | Benchmarks: Moody\u2019s A-rated health system medians')
    _add_footer(s2, 2)

    # ── SLIDE 3: STACKED AREA — Asset Composition ──
    s3 = prs.slides.add_slide(blank)
    _add_title(s3, 'Asset Composition Has Shifted Toward Investments and Away from Liquidity')

    cd3 = CategoryChartData()
    cd3.categories = YEARS
    cd3.add_series('Cash',          (693, 560, 2065, 1524, 846, 811, 835, 739))
    cd3.add_series('Patient AR',    (652, 692, 685, 897, 964, 930, 1049, 1246))
    cd3.add_series('LT Investments',(1077, 1210, 1215, 1467, 1422, 1608, 1744, 2102))
    cd3.add_series('PP&E Net',      (1321, 1433, 1635, 1821, 2037, 2230, 2317, 2318))
    cd3.add_series('Other Assets',  (1569, 1542, 1589, 1932, 1724, 1724, 1887, 2031))

    chart_shape = s3.shapes.add_chart(
        XL_CHART_TYPE.AREA_STACKED, Inches(1.0), Inches(1.2), Inches(11.3), Inches(5.0), cd3)
    chart = chart_shape.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.TOP
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(9)
    chart.legend.font.name = FONT

    colors3 = [AC_BLUE, AC_AMBER, AC_TEAL, AC_DEEP, AC_GHOST]
    for i, c in enumerate(colors3):
        _color_series(chart.plots[0], i, c)

    _style_chart_axes(chart, y_max=10000)
    _add_source(s3, 'Source: Audited Financial Statements, FY2018\u20132025. FY2020 cash spike = CARES Act / Medicare accelerated payments.')
    _add_footer(s3, 3)

    # ── SLIDE 4: WATERFALL (stacked bar simulation) ──
    s4 = prs.slides.add_slide(blank)
    _add_title(s4, 'Net Asset Growth: $2.5B Built Primarily on Investment Returns')

    cd4 = CategoryChartData()
    labels = ['FY2018\nNet Assets', 'Cumulative\nOp. Income', 'Cumulative\nInv. Returns',
              'Other\nChanges', 'FY2025\nNet Assets']
    cd4.categories = labels
    cd4.add_series('Base',     (0, 1788, 3098, 4394, 0))
    cd4.add_series('Increase', (1788, 1310, 1296, 0, 4289))
    cd4.add_series('Decrease', (0, 0, 0, 105, 0))

    chart_shape4 = s4.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_STACKED, Inches(1.5), Inches(1.2), Inches(10.3), Inches(5.2), cd4)
    chart4 = chart_shape4.chart
    chart4.has_legend = False

    s0 = chart4.plots[0].series[0]
    s0.format.fill.background()
    s0.format.line.fill.background()
    _color_series(chart4.plots[0], 1, AC_GREEN)
    _color_series(chart4.plots[0], 2, AC_RED)

    for cat_idx in [0, 4]:
        try:
            pt = chart4.plots[0].series[1].points[cat_idx]
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = AC_DEEP
        except Exception:
            pass

    chart4.plots[0].series[1].has_data_labels = True
    dl = chart4.plots[0].series[1].data_labels
    dl.font.name = FONT
    dl.font.size = Pt(11)
    dl.font.bold = True
    dl.font.color.rgb = AC_DARK
    dl.number_format = '$#,##0"M"'
    dl.position = XL_LABEL_POSITION.OUTSIDE_END

    _style_chart_axes(chart4, y_max=5000)
    _add_source(s4, 'Source: Audited Financial Statements, FY2018\u20132025')
    _add_footer(s4, 4)

    # ── SLIDE 5: GAUGE DASHBOARD (horizontal bar alternative) ──
    s5 = prs.slides.add_slide(blank)
    _add_title(s5, 'Rating Agency Dashboard: Two Green, Two Amber, Three Red')

    metrics = [
        ('Op. Margin',   2.4,  3.5,   6, '%',  False),
        ('EBIDA Margin', 5.5,  9.0,  12, '%',  False),
        ('Days Cash',    122,  225,  300, 'd',  False),
        ('Days AR',      51,   48,   70, 'd',  True),
        ('Debt/Cap',     28.5, 33,   55, '%',  True),
        ('CF Margin',    2.0,  7.0,  10, '%',  False),
        ('CapEx/Depr',   8,    110, 130, '%',  False),
    ]
    bar_top_start = Inches(1.3)
    bar_height = Inches(0.55)
    bar_gap = Inches(0.08)
    label_w = Inches(1.8)
    bar_w = Inches(7.5)
    val_w = Inches(2.0)

    for i, (label, actual, benchmark, mx, unit, invert) in enumerate(metrics):
        y = bar_top_start + Inches(i * 0.72)
        is_good = (actual <= benchmark) if invert else (actual >= benchmark)
        is_close = (actual <= benchmark * 1.2) if invert else (actual >= benchmark * 0.8)
        bar_color = AC_GREEN if is_good else (AC_AMBER if is_close else AC_RED)

        _add_textbox(s5, MARGIN_L, y + Pt(6), label_w, bar_height,
                     label, size=11, bold=True, color=AC_DARK)

        bg_bar = _add_rect(s5, MARGIN_L + label_w, y + Pt(4), bar_w, Inches(0.32), AC_LGRAY)
        fill_pct = min(actual / mx, 1.0)
        fill_w = Inches(7.5 * fill_pct)
        _add_rect(s5, MARGIN_L + label_w, y + Pt(4), fill_w, Inches(0.32), bar_color)

        bm_pct = min(benchmark / mx, 1.0)
        bm_x = MARGIN_L + label_w + Inches(7.5 * bm_pct)
        bm_line = s5.shapes.add_shape(MSO_SHAPE.RECTANGLE, bm_x, y + Pt(2), Pt(2), Inches(0.38))
        bm_line.fill.solid()
        bm_line.fill.fore_color.rgb = AC_AMBER
        bm_line.line.fill.background()

        val_text = f'{actual}{unit} vs {benchmark}{unit}'
        status_tag = '\u2713' if is_good else '\u26a0' if is_close else '\u2717'
        _add_textbox(s5, MARGIN_L + label_w + bar_w + Inches(0.15), y + Pt(6), val_w, bar_height,
                     f'{status_tag} {val_text}', size=9, bold=True,
                     color=bar_color, align=PP_ALIGN.LEFT)

    _add_source(s5, 'Source: Audited Financial Statements, FY2025 | Benchmarks: Moody\u2019s A-rated medians | Amber markers = benchmark target')
    _add_footer(s5, 5)

    # ── SLIDE 6: LEVERAGE — Combo bar+line ──
    s6 = prs.slides.add_slide(blank)
    _add_title(s6, 'Leverage Is Declining Rapidly \u2014 Driven by Net Asset Growth, Not Debt Paydown')

    cd6_bars = CategoryChartData()
    cd6_bars.categories = YEARS
    cd6_bars.add_series('Net Assets ($M)', (1788, 1796, 1700, 2748, 2692, 3127, 3737, 4289))

    chart_shape6 = s6.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1.0), Inches(1.2), Inches(11.3), Inches(5.0), cd6_bars)
    chart6 = chart_shape6.chart
    chart6.has_legend = True
    chart6.legend.position = XL_LEGEND_POSITION.TOP
    chart6.legend.include_in_layout = False
    chart6.legend.font.size = Pt(9)
    chart6.legend.font.name = FONT

    _color_series(chart6.plots[0], 0, AC_GHOST)

    cd6_line = CategoryChartData()
    cd6_line.categories = YEARS
    cd6_line.add_series('Debt/Cap %', (45.9, 46.7, 53.6, 38.0, 39.8, 36.3, 30.2, 28.5))

    chart6.add_chart(cd6_line, XL_CHART_TYPE.LINE)
    line_plot = chart6.plots[1]
    line_series = line_plot.series[0]
    line_series.format.line.color.rgb = AC_RED
    line_series.format.line.width = Pt(2.5)
    line_series.smooth = False
    line_series.has_data_labels = True
    dl6 = line_series.data_labels
    dl6.font.name = FONT
    dl6.font.size = Pt(9)
    dl6.font.bold = True
    dl6.font.color.rgb = AC_RED
    dl6.number_format = '0.0"%"'
    dl6.position = XL_LABEL_POSITION.ABOVE

    _style_chart_axes(chart6, y_max=5000)
    _add_source(s6, 'Source: Audited Financial Statements, FY2018\u20132025')
    _add_footer(s6, 6)

    # ── SLIDE 7: PENSION ──
    s7 = prs.slides.add_slide(blank)
    _add_title(s7, 'Pension Liability Virtually Eliminated: A $615M Turnaround Story')

    cd7 = CategoryChartData()
    cd7.categories = YEARS
    cd7.add_series('Pension Liability ($M)', (290, 437, 629, 305, 188, 94, 17, 14))

    chart_shape7 = s7.shapes.add_chart(
        XL_CHART_TYPE.AREA, Inches(1.0), Inches(1.2), Inches(11.3), Inches(3.5), cd7)
    chart7 = chart_shape7.chart
    chart7.has_legend = False
    _color_series(chart7.plots[0], 0, AC_RED, fill_opacity=15)
    chart7.plots[0].series[0].has_data_labels = True
    dl7 = chart7.plots[0].series[0].data_labels
    dl7.font.name = FONT
    dl7.font.size = Pt(10)
    dl7.font.bold = True
    dl7.font.color.rgb = AC_RED
    dl7.number_format = '$#,##0"M"'
    dl7.position = XL_LABEL_POSITION.ABOVE
    _style_chart_axes(chart7, y_max=700)

    _add_callout_box(s7, MARGIN_L, Inches(5.0), CONTENT_W, Inches(1.2),
                     'From $629M liability (FY2020) to $14M (FY2025) \u2014 now overfunded with $58.7M in plan assets',
                     'Plan appears frozen with improving funded status. Frees ~$30\u201350M in annual required contributions for strategic redeployment.',
                     AC_GREEN)
    _add_source(s7, 'Source: Audited Financial Statements + Note 6, FY2025')
    _add_footer(s7, 7)

    # ── SLIDE 8: LIQUIDITY — Stacked bar + line ──
    s8 = prs.slides.add_slide(blank)
    _add_title(s8, 'Liquidity: $2.8B in Investments, But Only 122 Days of Operating Cash')

    cd8_bars = CategoryChartData()
    cd8_bars.categories = YEARS
    cd8_bars.add_series('Cash ($M)',          (693, 560, 2065, 1524, 846, 811, 835, 739))
    cd8_bars.add_series('LT Investments ($M)',(1077, 1210, 1215, 1467, 1422, 1608, 1744, 2102))

    chart_shape8 = s8.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_STACKED, Inches(1.0), Inches(1.2), Inches(11.3), Inches(5.0), cd8_bars)
    chart8 = chart_shape8.chart
    chart8.has_legend = True
    chart8.legend.position = XL_LEGEND_POSITION.TOP
    chart8.legend.include_in_layout = False
    chart8.legend.font.size = Pt(9)
    chart8.legend.font.name = FONT

    _color_series(chart8.plots[0], 0, AC_BLUE)
    _color_series(chart8.plots[0], 1, AC_TEAL)

    cd8_line = CategoryChartData()
    cd8_line.categories = YEARS
    cd8_line.add_series('Days Cash (right)', (123, 121, 220, 175, 119, 120, 120, 122))
    chart8.add_chart(cd8_line, XL_CHART_TYPE.LINE)
    line_plot8 = chart8.plots[1]
    line_s8 = line_plot8.series[0]
    line_s8.format.line.color.rgb = AC_RED
    line_s8.format.line.width = Pt(2.5)
    line_s8.smooth = False
    line_s8.has_data_labels = True
    dl8 = line_s8.data_labels
    dl8.font.name = FONT
    dl8.font.size = Pt(9)
    dl8.font.bold = True
    dl8.font.color.rgb = AC_RED
    dl8.number_format = '0"d"'
    dl8.position = XL_LABEL_POSITION.ABOVE

    _style_chart_axes(chart8, y_max=4000)
    _add_source(s8, 'Source: Audited Financial Statements, FY2018\u20132025. FY2020 spike = CARES Act.')
    _add_footer(s8, 8)

    # ── SLIDE 9: DAYS AR ──
    s9 = prs.slides.add_slide(blank)
    _add_title(s9, 'Days in A/R Rising Steadily: 42 \u2192 51 Days, $204M in Consumed Working Capital')

    cd9 = CategoryChartData()
    cd9.categories = YEARS
    cd9.add_series('Days in A/R', (42.4, 44.4, 43.2, 48.7, 48.3, 43.9, 46.4, 50.7))

    chart_shape9 = s9.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1.0), Inches(1.2), Inches(11.3), Inches(3.3), cd9)
    chart9 = chart_shape9.chart
    chart9.has_legend = False

    _color_series(chart9.plots[0], 0, AC_BLUE)

    dar_vals = [42.4, 44.4, 43.2, 48.7, 48.3, 43.9, 46.4, 50.7]
    for i, v in enumerate(dar_vals):
        color = AC_GREEN if v < 45 else (AC_AMBER if v <= 48 else AC_RED)
        try:
            pt = chart9.plots[0].series[0].points[i]
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = color
        except Exception:
            pass

    chart9.plots[0].series[0].has_data_labels = True
    dl9 = chart9.plots[0].series[0].data_labels
    dl9.font.name = FONT
    dl9.font.size = Pt(10)
    dl9.font.bold = True
    dl9.font.color.rgb = AC_DARK
    dl9.number_format = '0.0"d"'
    dl9.position = XL_LABEL_POSITION.OUTSIDE_END
    _style_chart_axes(chart9, y_fmt='0"d"', y_min=35, y_max=55)

    _add_callout_box(s9, MARGIN_L, Inches(4.7), Inches(5.6), Inches(1.1),
                     '$204M in Working Capital Consumed',
                     '8.3-day deterioration \u00d7 $24.6M/day. AR grew 91% ($652M\u2192$1,246M) while revenue grew only 60%.',
                     AC_RED)
    _add_callout_box(s9, MARGIN_L + Inches(5.9), Inches(4.7), Inches(5.6), Inches(1.1),
                     'Commercial Payers Are the Friction Point',
                     '\u201cOther commercial\u201d = 26% of revenue but 37% of AR ($388M). Denial management and contract enforcement are the highest-yield targets.',
                     AC_AMBER)
    _add_source(s9, 'Source: Audited Financial Statements, FY2018\u20132025')
    _add_footer(s9, 9)

    # ── SLIDE 10: CAPEX vs DEPRECIATION ──
    s10 = prs.slides.add_slide(blank)
    _add_title(s10, 'Capital Reinvestment Crisis: CapEx Has Never Exceeded 31% of Depreciation')

    cd10 = CategoryChartData()
    cd10.categories = YEARS
    cd10.add_series('Depreciation ($M)', (206, 202, 205, 216, 219, 223, 248, 272))
    cd10.add_series('CapEx ($M)',        (21, 30, 33, 41, 51, 70, 22, 20))

    chart_shape10 = s10.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1.0), Inches(1.2), Inches(11.3), Inches(3.3), cd10)
    chart10 = chart_shape10.chart
    chart10.has_legend = True
    chart10.legend.position = XL_LEGEND_POSITION.TOP
    chart10.legend.include_in_layout = False
    chart10.legend.font.size = Pt(9)
    chart10.legend.font.name = FONT
    _color_series(chart10.plots[0], 0, AC_GHOST)
    _color_series(chart10.plots[0], 1, AC_RED)

    chart10.plots[0].series[1].has_data_labels = True
    dl10 = chart10.plots[0].series[1].data_labels
    dl10.font.name = FONT
    dl10.font.size = Pt(9)
    dl10.font.bold = True
    dl10.font.color.rgb = AC_RED
    ratios10 = ['10.1%','14.7%','16.2%','19.0%','23.4%','31.3%','9.1%','7.5%']
    dl10.number_format = '0.0"%"'
    dl10.position = XL_LABEL_POSITION.OUTSIDE_END
    _style_chart_axes(chart10, y_max=300)

    _add_callout_box(s10, MARGIN_L, Inches(4.7), CONTENT_W, Inches(1.2),
                     'Estimated $1.4B in Accumulated Deferred Maintenance (FY2018\u20132025)',
                     'With $2.3B in net PP&E and $272M in annual depreciation, the physical plant is aging without reinvestment. CapEx declined from $70M (FY2023) to $20M (FY2025) while revenue grew 16%.',
                     AC_RED)
    _add_source(s10, 'Source: Audited Financial Statements, FY2018\u20132025. Benchmark: CapEx/Depr = 100\u2013120%.')
    _add_footer(s10, 10)

    # ── SLIDE 11: SELF-INSURANCE ──
    s11 = prs.slides.add_slide(blank)
    _add_title(s11, 'Self-Insurance Reserves Are Growing: A Hidden Cash Drain')

    cd11 = CategoryChartData()
    cd11.categories = YEARS
    cd11.add_series('Self-Insurance ($M)', (283, 273, 282, 328, 359, 336, 363, 394))

    chart_shape11 = s11.shapes.add_chart(
        XL_CHART_TYPE.AREA, Inches(1.0), Inches(1.2), Inches(11.3), Inches(3.0), cd11)
    chart11 = chart_shape11.chart
    chart11.has_legend = False
    _color_series(chart11.plots[0], 0, AC_AMBER, fill_opacity=20)
    chart11.plots[0].series[0].has_data_labels = True
    dl11 = chart11.plots[0].series[0].data_labels
    dl11.font.name = FONT
    dl11.font.size = Pt(10)
    dl11.font.bold = True
    dl11.font.color.rgb = AC_DARK
    dl11.number_format = '$#,##0"M"'
    dl11.position = XL_LABEL_POSITION.ABOVE
    _style_chart_axes(chart11, y_min=200, y_max=420)

    _add_callout_box(s11, MARGIN_L, Inches(4.5), Inches(5.6), Inches(1.1),
                     '$394M in Reserves, Up 39% Since FY2018',
                     'Professional & general liability accruals: $448M (FY2025). Per-claim retention rose from $20M to $25M.',
                     AC_AMBER)
    _add_callout_box(s11, MARGIN_L + Inches(5.9), Inches(4.5), Inches(5.6), Inches(1.1),
                     'Why It Matters',
                     'Self-insurance is a long-tail liability. Today\u2019s reserves fund yesterday\u2019s claims. If actuarial estimates prove inadequate, future cash calls accelerate.',
                     AC_BLUE)
    _add_source(s11, 'Source: Audited Financial Statements + Note 11, FY2018\u20132025')
    _add_footer(s11, 11)

    # ── SLIDE 12: INVESTMENT RETURNS vs OP INCOME ──
    s12 = prs.slides.add_slide(blank)
    _add_title(s12, 'Investment Returns Are Masking Weak Operating Performance')

    cd12 = CategoryChartData()
    cd12.categories = YEARS
    cd12.add_series('Operating Income ($M)',   (161, 152, 131, 260, 81, 147, 159, 219))
    cd12.add_series('Investment Returns ($M)', (119, 53, 20, 511, -271, 214, 359, 346))

    chart_shape12 = s12.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1.0), Inches(1.2), Inches(11.3), Inches(5.0), cd12)
    chart12 = chart_shape12.chart
    chart12.has_legend = True
    chart12.legend.position = XL_LEGEND_POSITION.TOP
    chart12.legend.include_in_layout = False
    chart12.legend.font.size = Pt(9)
    chart12.legend.font.name = FONT
    _color_series(chart12.plots[0], 0, AC_BLUE)
    _color_series(chart12.plots[0], 1, AC_TEAL)

    try:
        pt = chart12.plots[0].series[1].points[4]
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = AC_RED
    except Exception:
        pass

    _style_chart_axes(chart12, y_min=-350, y_max=550)

    _add_textbox(s12, Inches(8.5), Inches(1.3), Inches(4), Inches(0.3),
                 'Cumulative: Op Inc $1,310M vs Inv Returns $1,351M',
                 size=10, bold=True, color=AC_DARK, align=PP_ALIGN.RIGHT)
    _add_source(s12, 'Source: Audited Financial Statements, FY2018\u20132025')
    _add_footer(s12, 12)

    # ── SLIDE 13: OPPORTUNITY TABLE ──
    s13 = prs.slides.add_slide(blank)
    _add_title(s13, 'Balance Sheet Opportunities: Protect Strengths, Close the Gaps')

    from pptx.util import Inches as In
    rows_data = [
        ('Domain', 'Status', 'Action', 'Impact'),
        ('Debt / Capitalization', '\u2713 STRONG', 'Protect. Use 450bp headroom as strategic debt capacity ($300\u2013400M).', 'Enabling'),
        ('Pension Elimination', '\u2713 STRONG', 'Redirect $30\u201350M in freed annual contributions to capital investment.', '$30\u201350M/yr'),
        ('Days in A/R', '\u26a0 WARNING', 'Reduce from 51d to 45d via commercial denial management. One-time $148M cash release.', '$148M + ongoing'),
        ('Days Cash on Hand', '\u2717 BELOW', 'Improve CF via margin gains, release AR working capital, rebalance investment liquidity.', '+78 days target'),
        ('CapEx / Depreciation', '\u2717 CRITICAL', 'Facilities master plan. Phase $1.4B deferred backlog over 5\u20137 years.', '$1.0\u20131.4B plan'),
        ('Self-Insurance', '\u26a0 RISING', 'Actuarial adequacy review. Explore risk transfer for tail exposure.', 'Risk mitigation'),
    ]
    col_widths = [In(2.0), In(1.3), In(6.5), In(1.9)]
    tbl = s13.shapes.add_table(len(rows_data), 4, MARGIN_L, In(1.2), In(11.7), In(4.5)).table

    for ci, w in enumerate(col_widths):
        tbl.columns[ci].width = w

    status_colors = {'\u2713': AC_GREEN, '\u26a0': AC_AMBER, '\u2717': AC_RED}
    for ri, row in enumerate(rows_data):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = ''
            p = cell.text_frame.paragraphs[0]
            if ri == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = AC_BLUE
                _add_run(p, val, size=9, bold=True, color=AC_WHITE)
                p.alignment = PP_ALIGN.CENTER
            else:
                if ri % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = AC_ALT
                sc = None
                if ci == 1:
                    for k, v in status_colors.items():
                        if val.startswith(k):
                            sc = v
                            break
                _add_run(p, val, size=9, bold=(ci <= 1), color=sc if sc else AC_DARK)

            cell.margin_left = Pt(6)
            cell.margin_right = Pt(6)
            cell.margin_top = Pt(4)
            cell.margin_bottom = Pt(4)

    for ri in range(len(rows_data)):
        for ci in range(4):
            cell = tbl.cell(ri, ci)
            for edge in ['top', 'bottom', 'left', 'right']:
                border = cell._tc.get_or_add_tcPr()
                bdr_el = border.makeelement(qn(f'a:ln{edge.capitalize()[0:1]}'), {})
                if ri == 0 or ri == len(rows_data) - 1 or ci == 0 or ci == 3:
                    bdr_el.set('w', str(Pt(1)))
                else:
                    bdr_el.set('w', str(Pt(0.5)))

    _add_source(s13, 'Source: Analysis based on Audited Financial Statements, FY2018\u20132025')
    _add_footer(s13, 13)

    # ── SLIDE 14: ROADMAP ──
    s14 = prs.slides.add_slide(blank)
    _add_title(s14, 'Recommended Engagement: Protect, Release, Rebuild')

    phases = [
        (AC_BLUE, 'Phase 1: Protect (0\u20136 mo)',
         '\u2022 Revenue cycle diagnostic \u2014 commercial AR deep dive\n\u2022 Working capital optimization \u2014 AR to 45 days\n\u2022 Self-insurance actuarial adequacy review\n\u2022 Investment portfolio liquidity assessment'),
        (AC_DEEP, 'Phase 2: Release (6\u201318 mo)',
         '\u2022 Freed pension contributions \u2192 capital reserve\n\u2022 AR cash release \u2192 liquidity cushion\n\u2022 Operating margin improvement \u2192 CF generation\n\u2022 Debt capacity sizing for strategic deployment'),
        (AC_MID, 'Phase 3: Rebuild (18\u201336 mo)',
         '\u2022 Facilities master plan & capital allocation\n\u2022 Phased deferred maintenance program\n\u2022 Rating agency engagement strategy\n\u2022 Long-term investment policy rebalancing'),
    ]
    phase_w = Inches(3.7)
    for i, (bg, title, bullets) in enumerate(phases):
        left = MARGIN_L + Inches(i * 3.85)
        top = Inches(1.3)
        box = _add_rect(s14, left, top, phase_w, Inches(3.2), bg)
        tb = s14.shapes.add_textbox(left + Pt(12), top + Pt(10), phase_w - Pt(24), Inches(3.0))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        _add_run(p, title, size=12, bold=True, color=AC_WHITE)
        for line in bullets.split('\n'):
            p2 = tf.add_paragraph()
            p2.space_before = Pt(3)
            _add_run(p2, line, size=9, color=AC_WHITE)

    _add_callout_box(s14, MARGIN_L, Inches(4.8), CONTENT_W, Inches(1.1),
                     'MedStar has the balance sheet to fund its own transformation',
                     '$300\u2013400M in debt capacity headroom + $30\u201350M/yr in freed pension contributions + $148M in AR cash release = the capital exists. What\u2019s needed is a disciplined framework to deploy it.',
                     AC_GREEN)
    _add_footer(s14, 14)

    # ── SLIDE 15: APPENDIX TABLE ──
    s15 = prs.slides.add_slide(blank)
    _add_title(s15, 'Appendix: Key Balance Sheet Ratios (FY2018\u20132025)')

    hdrs = ['Metric','FY18','FY19','FY20','FY21','FY22','FY23','FY24','FY25','A-Rated']
    data_rows = [
        ['Days Cash','123','121','220','175','119','120','120','122','225'],
        ['Days AR','42','44','43','49','48','44','46','51','48'],
        ['Debt/Cap %','45.9','46.7','53.6','38.0','39.8','36.3','30.2','28.5','33'],
        ['Debt/EBIDA','4.6x','4.6x','6.0x','3.7x','6.1x','5.0x','4.1x','3.6x','3.0x'],
        ['CapEx/Depr %','10','15','16','19','23','31','9','8','110'],
        ['CF Margin %','4.0','4.8','18.4','5.7','\u22126.0','3.6','3.7','2.0','7.0'],
        ['Pension ($M)','290','437','629','305','188','94','17','14','\u2014'],
        ['Self-Ins ($M)','283','273','282','328','359','336','363','394','\u2014'],
        ['Net Assets ($M)','1,788','1,796','1,700','2,748','2,692','3,127','3,737','4,289','\u2014'],
    ]
    all_rows = [hdrs] + data_rows
    tbl15 = s15.shapes.add_table(len(all_rows), 10, MARGIN_L, Inches(1.2), Inches(11.7), Inches(5.0)).table

    col_ws = [Inches(1.6)] + [Inches(1.01)] * 8 + [Inches(1.1)]
    for ci, w in enumerate(col_ws):
        tbl15.columns[ci].width = w

    for ri, row in enumerate(all_rows):
        for ci, val in enumerate(row):
            cell = tbl15.cell(ri, ci)
            cell.text = ''
            p = cell.text_frame.paragraphs[0]
            if ri == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = AC_BLUE
                _add_run(p, val, size=8, bold=True, color=AC_WHITE)
                p.alignment = PP_ALIGN.CENTER
            elif ci == 9:
                cell.fill.solid()
                cell.fill.fore_color.rgb = AC_AMBER
                _add_run(p, val, size=8, bold=True, color=AC_WHITE)
                p.alignment = PP_ALIGN.RIGHT
            else:
                if ri % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = AC_ALT
                _add_run(p, val, size=8, bold=(ci == 0), color=AC_DARK)
                p.alignment = PP_ALIGN.RIGHT if ci > 0 else PP_ALIGN.LEFT

            cell.margin_left = Pt(4)
            cell.margin_right = Pt(4)
            cell.margin_top = Pt(3)
            cell.margin_bottom = Pt(3)

    _add_source(s15, 'Source: Audited Financial Statements, FY2018\u20132025 | Benchmarks: Moody\u2019s A-rated medians (2024)')
    _add_footer(s15, 15)

    prs.save(OUT_PATH)
    print(f'Saved: {OUT_PATH}')


if __name__ == '__main__':
    build_deck()
